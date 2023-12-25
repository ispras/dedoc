from typing import Dict, List, Optional

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.picture import Picture
from pptx.slide import Slide

from dedoc.attachments_extractors.concrete_attachments_extractors.pptx_attachments_extractor import PptxAttachmentsExtractor
from dedoc.data_structures import AttachAnnotation, Table, TableAnnotation
from dedoc.data_structures.cell_with_meta import CellWithMeta
from dedoc.data_structures.line_metadata import LineMetadata
from dedoc.data_structures.line_with_meta import LineWithMeta
from dedoc.data_structures.table_metadata import TableMetadata
from dedoc.data_structures.unstructured_document import UnstructuredDocument
from dedoc.extensions import recognized_extensions, recognized_mimes
from dedoc.readers.base_reader import BaseReader
from dedoc.utils.utils import get_mime_extension


class PptxReader(BaseReader):
    """
    This class is used for parsing documents with .pptx extension.
    Please use :class:`~dedoc.converters.PptxConverter` for getting pptx file from similar formats.
    """

    def __init__(self, *, config: Optional[dict] = None) -> None:
        super().__init__(config=config)
        self.attachments_extractor = PptxAttachmentsExtractor(config=self.config)

    def can_read(self, file_path: Optional[str] = None, mime: Optional[str] = None, extension: Optional[str] = None, parameters: Optional[dict] = None) -> bool:
        """
        Check if the document extension is suitable for this reader.
        Look to the documentation of :meth:`~dedoc.readers.BaseReader.can_read` to get information about the method's parameters.
        """
        extension, mime = get_mime_extension(file_path=file_path, mime=mime, extension=extension)
        return extension.lower() in recognized_extensions.pptx_like_format or mime in recognized_mimes.pptx_like_format

    def read(self, file_path: str, parameters: Optional[dict] = None) -> UnstructuredDocument:
        """
        The method return document content with all document's lines, tables and attachments.
        Look to the documentation of :meth:`~dedoc.readers.BaseReader.read` to get information about the method's parameters.
        """
        parameters = {} if parameters is None else parameters

        with_attachments = self.attachments_extractor.with_attachments(parameters=parameters)
        attachments = self.attachments_extractor.extract(file_path=file_path, parameters=parameters) if with_attachments else []
        attachment_name2uid = {attachment.original_name: attachment.uid for attachment in attachments}

        prs = Presentation(file_path)
        lines, tables = [], []

        for page_id, slide in enumerate(prs.slides, start=1):
            images_rels = self.__get_slide_images_rels(slide)

            for paragraph_id, shape in enumerate(slide.shapes, start=1):

                if shape.has_text_frame:
                    lines.append(LineWithMeta(line=f"{shape.text}\n", metadata=LineMetadata(page_id=page_id, line_id=paragraph_id)))

                if shape.has_table:
                    self.__add_table(lines, tables, page_id, paragraph_id, shape)

                if with_attachments and hasattr(shape, "image"):
                    if len(lines) == 0:
                        lines.append(LineWithMeta(line="", metadata=LineMetadata(page_id=page_id, line_id=paragraph_id)))
                    self.__add_attach_annotation(lines[-1], shape, attachment_name2uid, images_rels)

        return UnstructuredDocument(lines=lines, tables=tables, attachments=attachments, warnings=[])

    def __add_table(self, lines: List[LineWithMeta], tables: List[Table], page_id: int, paragraph_id: int, shape: GraphicFrame) -> None:
        cells = [
            [CellWithMeta(lines=[LineWithMeta(line=cell.text, metadata=LineMetadata(page_id=page_id, line_id=0))]) for cell in row.cells]
            for row in shape.table.rows
        ]
        table = Table(cells=cells, metadata=TableMetadata(page_id=page_id))

        if len(lines) == 0:
            lines.append(LineWithMeta(line="", metadata=LineMetadata(page_id=page_id, line_id=paragraph_id)))
        lines[-1].annotations.append(TableAnnotation(start=0, end=len(lines[-1]), name=table.metadata.uid))
        tables.append(table)

    def __get_slide_images_rels(self, slide: Slide) -> Dict[str, str]:
        rels = BeautifulSoup(slide.part.rels.xml, "xml")
        images_dir = "../media/"

        images_rels = dict()
        for rel in rels.find_all("Relationship"):
            if rel["Target"].startswith(images_dir):
                images_rels[rel["Id"]] = rel["Target"][len(images_dir):]

        return images_rels

    def __add_attach_annotation(self, line: LineWithMeta, shape: Picture, attachment_name2uid: dict, images_rels: dict) -> None:
        try:
            image_rels_id = shape.element.blip_rId
            image_name = images_rels[image_rels_id]
            image_uid = attachment_name2uid[image_name]
            line.annotations.append(AttachAnnotation(start=0, end=len(line), attach_uid=image_uid))
        except KeyError as e:
            self.logger.warning(f"Attachment key hasn't been found ({e})")
